"""Build the real 36x24 poster PPTX from scratch.

Layout follows poster_mockup_v4.svg, minus the bottom takeaway strip
(dropped per user request). Produces poster/poster_final.pptx.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

POSTER = Path(__file__).parent
REPORT_FIG = POSTER.parent / "report" / "figures"
POSTER_FIG = POSTER / "figures"

# -- palette ---------------------------------------------------------
RED = RGBColor(0xCC, 0x06, 0x33)          # SFU red
RED_SOFT = RGBColor(0xF8, 0xF0, 0xEF)     # motivation band bg
OUR_SOFT = RGBColor(0xFF, 0xF0, 0xEC)     # LCM row tint
OUR_STRONG = RGBColor(0xFF, 0xE4, 0xDC)   # BPN + Refiner row tint
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)
GREY = RGBColor(0xBB, 0xBB, 0xBB)
POS = RGBColor(0x20, 0x70, 0x30)
NEG = RGBColor(0xA0, 0x20, 0x20)


def add_textbox(slide, x, y, w, h, text, *, size=16, bold=False,
                italic=False, color=INK, align=PP_ALIGN.LEFT,
                font="Calibri", anchor=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor

    # Build runs: pass list of (text, style_dict) OR plain str.
    if isinstance(text, str):
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
            r.font.name = font
    else:
        # list of paragraphs, each paragraph is list of runs: [(text, {"bold":..}), ...]
        for i, para in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            for (run_text, run_style) in para:
                r = p.add_run()
                r.text = run_text
                r.font.size = Pt(run_style.get("size", size))
                r.font.bold = run_style.get("bold", bold)
                r.font.italic = run_style.get("italic", italic)
                r.font.color.rgb = run_style.get("color", color)
                r.font.name = run_style.get("font", font)
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_pill(slide, x, y, w, h, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = RED
    shape.line.width = Pt(2.5)
    tf = shape.text_frame
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = RED
    r.font.name = "Calibri"


def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                    Inches(w), Inches(h))


# -- build -----------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(36)
prs.slide_height = Inches(24)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ========== TITLE BAR (y=0..3.0) ==========
add_rect(slide, 0, 0, 36, 3.0, RED)

# Title, subtitle, authors (centered)
add_textbox(slide, 4.0, 0.15, 28.0, 1.1,
            "Cross-Language Scene Text Replacement in Video",
            size=64, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=None)
add_textbox(slide, 4.0, 1.25, 28.0, 0.7,
            "A Modular Five-Stage Pipeline",
            size=36, italic=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER)
add_textbox(slide, 4.0, 2.0, 28.0, 0.6,
            "Hebin Yao  ·  Yunshan Feng  ·  Liliana Lopez  "
            "·  SFU School of Computing Science  ·  CMPT 743 Visual Computing Lab II",
            size=24, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER)

# Left logo placeholder
add_rect(slide, 0.4, 0.3, 3.2, 2.4, RGBColor(0xFF, 0xFF, 0xFF),
         line=RGBColor(0xFF, 0xFF, 0xFF), line_width=0.5)
add_rect(slide, 0.4, 0.3, 3.2, 2.4, RED)  # keep background red
add_textbox(slide, 0.4, 1.2, 3.2, 0.6, "logos here",
            size=18, italic=True, color=RGBColor(0xFF, 0xCC, 0xD6),
            align=PP_ALIGN.CENTER)

# Right QR placeholder
add_rect(slide, 32.4, 0.3, 3.2, 2.4, RGBColor(0xFF, 0xFF, 0xFF),
         line=RGBColor(0xFF, 0xFF, 0xFF), line_width=0.5)
add_textbox(slide, 32.4, 1.0, 3.2, 0.5, "GitHub QR",
            size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_textbox(slide, 32.4, 1.5, 3.2, 0.5, "(scan for code + demo)",
            size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# ========== MOTIVATION BAND (y=3.15..6.95, x=0.4..25.3) ==========
MOT_X, MOT_Y, MOT_W, MOT_H = 0.4, 3.15, 24.9, 4.05
add_rect(slide, MOT_X, MOT_Y, MOT_W, MOT_H, RED_SOFT, line=RED, line_width=3)

# headline
add_textbox(slide, MOT_X + 0.25, MOT_Y + 0.1, MOT_W - 0.5, 0.7,
            "Replacing scene text inside video, across languages",
            size=38, bold=True, color=RED)

# problem 3-line
add_textbox(slide, MOT_X + 0.25, MOT_Y + 0.85, MOT_W - 0.5, 1.1,
            [
                [
                    ("Video localization today means dubbing and subtitling — but ", {"size": 18}),
                    ("scene text inside the frame", {"size": 18, "bold": True}),
                    (" (signs, menus, banners, product", {"size": 18}),
                ],
                [("labels) stays in the source language. Manual rotoscoping and 3-D tracking are the current fallback; neither scales.",
                  {"size": 18})],
                [("A usable automated replacement must satisfy four axes simultaneously:",
                  {"size": 18})],
            ])

# 4-axes pill row
pill_w = (MOT_W - 0.5 - 0.3 * 3) / 4
pill_y = MOT_Y + 2.0
labels = ["(i) Translated", "(ii) Aligned",
          "(iii) Temporally consistent", "(iv) Photorealistic"]
for i, lbl in enumerate(labels):
    add_pill(slide, MOT_X + 0.25 + i * (pill_w + 0.3),
             pill_y, pill_w, 0.45, lbl)

# Prior-work 2-col strip
pw_y = MOT_Y + 2.60
card_w = (MOT_W - 0.5 - 0.3) / 2
# left card: STRIVE
add_rect(slide, MOT_X + 0.25, pw_y, card_w, 1.0,
         RGBColor(0xFF, 0xFF, 0xFF), line=GREY, line_width=1)
add_textbox(slide, MOT_X + 0.35, pw_y + 0.05, card_w - 0.2, 0.9,
            [
                [("Closest prior work — STRIVE (ICCV 2021)",
                  {"size": 18, "bold": True, "color": INK})],
                [("+ frontalize → edit → propagate → composite structure for video text replacement",
                  {"size": 16, "color": POS})],
                [("− code unreleased; demos only in-language edits, not translation",
                  {"size": 16, "color": NEG})],
            ])
# right card: generative editors
add_rect(slide, MOT_X + 0.25 + card_w + 0.3, pw_y, card_w, 1.0,
         RGBColor(0xFF, 0xFF, 0xFF), line=GREY, line_width=1)
add_textbox(slide, MOT_X + 0.35 + card_w + 0.3, pw_y + 0.05,
            card_w - 0.2, 0.9,
            [
                [("Different paradigm — generative editors (Sora, Gen-4 Aleph)",
                  {"size": 18, "bold": True, "color": INK})],
                [("+ naturalistic style, lighting, and motion",
                  {"size": 16, "color": POS})],
                [("− cannot reliably render a specified target string",
                  {"size": 16, "color": NEG})],
            ])

# bridging tagline
add_textbox(slide, MOT_X + 0.25, pw_y + 1.05, MOT_W - 0.5, 0.35,
            "Our response — a reproducible, open, cross-language pipeline "
            "(full details below: diagram, custom modules, contribution table).",
            size=18, bold=True, italic=True, color=RED,
            align=PP_ALIGN.CENTER)

# ========== LEFT COLUMN — 5-Stage Pipeline (x=0.4..5.5) ==========
LX, LY, LW = 0.4, 7.40, 5.1
add_rect(slide, LX, LY, LW, 12.40, RGBColor(0xFA, 0xFA, 0xFA),
         line=GREY, line_width=1)
add_textbox(slide, LX + 0.2, LY + 0.15, LW - 0.4, 0.6,
            "The 5-Stage Pipeline",
            size=26, bold=True, color=RED)

stages = [
    ("S1 · Detect & Track",
     ["PaddleOCR quads · wordfreq gibberish filter",
      "CoTracker3 propagates quads (60-frame windows)",
      "4-filter funnel selects reference frame"]),
    ("S2 · Frontalize + Align  (ours)",
     ["RANSAC homography → canonical rectangle",
      "Alignment Refiner collapses ΔH to sub-pixel"]),
    ("S3 · Edit",
     ["AnyText2 (ICLR '25) writes target string",
      "Adaptive mask fixes long→short translation"]),
    ("S4 · Propagate  (BPN ours)",
     ["LCM — per-pixel lighting ratio map (re-impl)",
      "BPN — per-frame blur (re-impl, trained on our data)"]),
    ("S5 · Revert & Composite",
     ["Inverse warp + Poisson blend"]),
]
ty = LY + 0.85
for title, bullets in stages:
    add_textbox(slide, LX + 0.2, ty, LW - 0.4, 0.45,
                title, size=18, bold=True, color=RED)
    for i, b in enumerate(bullets):
        add_textbox(slide, LX + 0.25, ty + 0.5 + i * 0.38, LW - 0.4, 0.4,
                    "• " + b, size=14, color=INK)
    ty += 0.5 + 0.38 * len(bullets) + 0.35

add_textbox(slide, LX + 0.2, LY + 11.80, LW - 0.4, 0.5,
            "Threaded by one TextTrack dataclass (report §3.2).",
            size=13, italic=True, color=MUTED)

# ========== CENTER · UPPER — Pipeline diagram (x=5.7..20.0) ==========
CX, CY, CW, CH = 5.7, 7.40, 14.3, 8.0
add_rect(slide, CX, CY, CW, CH, RGBColor(0xFF, 0xFF, 0xFF),
         line=RED, line_width=3)
add_textbox(slide, CX, CY + 0.08, CW, 0.5,
            "Pipeline", size=26, bold=True, color=RED,
            align=PP_ALIGN.CENTER)

pipe_img = POSTER_FIG / "pipeline.png"
if pipe_img.exists():
    add_image(slide, pipe_img, CX + 0.2, CY + 0.65, CW - 0.4, CH - 1.15)
add_textbox(slide, CX, CY + CH - 0.45, CW, 0.4,
            "Five stages, one TextTrack dataclass — Detect → Frontalize → Edit → Propagate → Revert",
            size=16, italic=True, color=INK, align=PP_ALIGN.CENTER)

# ========== CENTER · LOWER — Our contribution table (x=5.7..20.0) ==========
TX, TY, TW, TH = 5.7, 15.55, 14.3, 4.25
add_rect(slide, TX, TY, TW, TH, RGBColor(0xFA, 0xFA, 0xFA),
         line=RED, line_width=3)

# heading
add_textbox(slide, TX + 0.2, TY + 0.08, TW - 0.4, 0.5,
            [[
                ("Our contribution — ", {"size": 24, "bold": True, "color": RED}),
                ("7 integrated models, 2 trained by us",
                 {"size": 24, "bold": True, "color": RED}),
            ]])
add_textbox(slide, TX + 0.2, TY + 0.55, TW - 0.4, 0.35,
            "A reproducible, open, cross-language baseline. Our engineering concentrates on the bottom two rows; the other five are integrated as-is.",
            size=13, italic=True, color=MUTED)

# table
rows = [
    ("Model", "Stage", "Our involvement", "header"),
    ("PaddleOCR (PP-OCRv4)", "S1", "Pretrained, used as is", "even"),
    ("CoTracker3", "S1", "Pretrained, used as is", "odd"),
    ("AnyText2 (ICLR 2025)", "S3", "Pretrained (Gradio server)", "even"),
    ("SRNet", "S3/4/5", "Pretrained, used as is", "odd"),
    ("Hi-SAM", "S3/4/5", "Pretrained, default inpainter", "even"),
    ("LCM", "S4", "Re-implementation of STRIVE's algorithm", "lcm"),
    ("BPN", "S4", "Re-impl of STRIVE · trained on our S2-aligned data", "ours"),
    ("Alignment Refiner", "S2", "Designed & trained from scratch", "ours"),
]

row_h = 0.33
tbl_x = TX + 0.2
tbl_y = TY + 1.0
tbl_w = TW - 0.4
col_w = [5.3, 1.6, tbl_w - 5.3 - 1.6]  # model, stage, involvement

for i, (m, s, inv, kind) in enumerate(rows):
    y = tbl_y + i * row_h
    if kind == "header":
        add_rect(slide, tbl_x, y, tbl_w, row_h, RED)
        text_color = RGBColor(0xFF, 0xFF, 0xFF)
        bold = True
    elif kind == "odd":
        add_rect(slide, tbl_x, y, tbl_w, row_h,
                 RGBColor(0xFA, 0xFA, 0xFA),
                 line=RGBColor(0xE8, 0xE8, 0xE8), line_width=0.5)
        text_color = INK
        bold = False
    elif kind == "even":
        add_rect(slide, tbl_x, y, tbl_w, row_h,
                 RGBColor(0xFF, 0xFF, 0xFF),
                 line=RGBColor(0xE8, 0xE8, 0xE8), line_width=0.5)
        text_color = INK
        bold = False
    elif kind == "lcm":
        add_rect(slide, tbl_x, y, tbl_w, row_h, OUR_SOFT,
                 line=RGBColor(0xE8, 0xE8, 0xE8), line_width=0.5)
        text_color = INK
        bold = False
    elif kind == "ours":
        add_rect(slide, tbl_x, y, tbl_w, row_h, OUR_STRONG,
                 line=RGBColor(0xE8, 0xE8, 0xE8), line_width=0.5)
        text_color = INK
        bold = True

    cx = tbl_x
    for j, (val, cw) in enumerate(zip([m, s, inv], col_w)):
        color = text_color
        if kind == "ours" and j == 2:
            color = RED
        if kind == "lcm" and j == 2:
            color = RED
        add_textbox(slide, cx + 0.1, y + 0.03, cw - 0.2, row_h - 0.04,
                    val, size=15, bold=bold, color=color)
        cx += cw

# recap footer
add_textbox(slide, TX + 0.2, tbl_y + len(rows) * row_h + 0.08,
            TW - 0.4, 0.3,
            "→ Alignment Refiner (S2): geometric alignment.    "
            "BPN (S4): per-frame blur adaptation.",
            size=13, italic=True, color=MUTED)

# ========== CUSTOM MODULES column (x=20.2..25.3) ==========
MX, MY, MW = 20.2, 7.40, 5.1
add_rect(slide, MX, MY, MW, 12.40, RGBColor(0xFA, 0xFA, 0xFA),
         line=GREY, line_width=1)
add_textbox(slide, MX + 0.2, MY + 0.15, MW - 0.4, 0.9,
            "Custom modules for video text editing",
            size=22, bold=True, color=RED)

# Alignment Refiner
sub_y = MY + 1.15
add_textbox(slide, MX + 0.2, sub_y, MW - 0.4, 0.4,
            "Alignment Refiner", size=18, bold=True, color=RED)
add_textbox(slide, MX + 0.2, sub_y + 0.4, MW - 0.4, 0.35,
            "Homography network for precise ROI tracking",
            size=13, italic=True, color=MUTED)
refiner_img = REPORT_FIG / "refiner-network.png"
if refiner_img.exists():
    add_image(slide, refiner_img, MX + 0.3, sub_y + 0.85,
              MW - 0.6, 1.8)

# LCM
sub_y2 = sub_y + 3.0
add_textbox(slide, MX + 0.2, sub_y2, MW - 0.4, 0.4,
            "LCM (Lighting Correction)", size=18, bold=True, color=RED)
add_textbox(slide, MX + 0.2, sub_y2 + 0.4, MW - 0.4, 0.35,
            "Per-frame lighting synchronization",
            size=13, italic=True, color=MUTED)
# Use lcm-result.jpg (SAMPLE→MUESTRA) — crop-preserving centered
lcm_mod_img = REPORT_FIG / "lcm-result.jpg"
if lcm_mod_img.exists():
    # Tall aspect: use a narrower centered slot
    add_image(slide, lcm_mod_img, MX + 1.3, sub_y2 + 0.85,
              MW - 2.6, 1.8)

# BPN
sub_y3 = sub_y2 + 3.0
add_textbox(slide, MX + 0.2, sub_y3, MW - 0.4, 0.4,
            "BPN (Blur Prediction)", size=18, bold=True, color=RED)
add_textbox(slide, MX + 0.2, sub_y3 + 0.4, MW - 0.4, 0.35,
            "Blur-kernel prediction for text consistency",
            size=13, italic=True, color=MUTED)
bpn_img = REPORT_FIG / "bpn.png"
if bpn_img.exists():
    add_image(slide, bpn_img, MX + 0.3, sub_y3 + 0.85,
              MW - 0.6, 1.8)

# ========== RIGHT COLUMN — RESULTS (x=25.5..35.6) ==========
RX, RY, RW, RH = 25.5, 3.15, 10.1, 19.25
add_rect(slide, RX, RY, RW, RH, RGBColor(0xFA, 0xFA, 0xFA),
         line=RED, line_width=3)

add_textbox(slide, RX + 0.25, RY + 0.15, RW - 0.5, 0.55,
            "Results", size=32, bold=True, color=RED)
add_textbox(slide, RX + 0.25, RY + 0.75, RW - 0.5, 0.9,
            "End-to-end pipeline outputs. The specified target string is "
            "rendered in place; perspective, lighting, and motion blur are "
            "inherited from the surrounding scene.",
            size=15, italic=True, color=MUTED)

# PS5 (top)
ps5_img = REPORT_FIG / "results_overview.png"
if ps5_img.exists():
    add_image(slide, ps5_img, RX + 0.2, RY + 1.75, RW - 0.4, 5.2)
add_textbox(slide, RX + 0.25, RY + 7.0, RW - 0.5, 0.6,
            "End-to-end: English → Simplified Chinese. Top = source frames, "
            "bottom = pipeline outputs (report Fig. 7).",
            size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# divider
add_rect(slide, RX + 0.25, RY + 7.72, RW - 0.5, 0.025, RED)

# SHIRTS
add_textbox(slide, RX + 0.25, RY + 7.85, RW - 0.5, 0.45,
            "Inpainter backends (S3)", size=20, bold=True, color=RED)
shirts_img = REPORT_FIG / "inpainting.png"
if shirts_img.exists():
    add_image(slide, shirts_img, RX + 0.2, RY + 8.3, RW - 0.4, 2.8)
add_textbox(slide, RX + 0.25, RY + 11.15, RW - 0.5, 0.5,
            "Source \"SHIRTS\" → background under SRNet (top) vs. Hi-SAM (bottom).",
            size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# divider
add_rect(slide, RX + 0.25, RY + 11.78, RW - 0.5, 0.025, RED)

# SAMPLE TEXT / LCM
add_textbox(slide, RX + 0.25, RY + 11.9, RW - 0.5, 0.45,
            "Lighting preservation (LCM)", size=20, bold=True, color=RED)
lcm_img = REPORT_FIG / "lcm-result.jpg"
if lcm_img.exists():
    add_image(slide, lcm_img, RX + 1.5, RY + 12.35, RW - 3.0, 6.25)
add_textbox(slide, RX + 0.25, RY + 18.7, RW - 0.5, 0.5,
            "SAMPLE TEXT → MUESTRA TEXTO across 8 lighting conditions.",
            size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# ========== FOOTER (y=22.5..24) ==========
add_rect(slide, 0, 22.55, 36, 1.45, RGBColor(0xFF, 0xFF, 0xFF))
add_rect(slide, 0, 22.55, 36, 0.04, RED)

# SFU placeholder logo (left)
add_rect(slide, 0.5, 22.75, 3.5, 1.05, RGBColor(0xFF, 0xFF, 0xFF),
         line=GREY, line_width=1)
add_textbox(slide, 0.5, 23.05, 3.5, 0.5,
            "SFU School of Computing Science",
            size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# center tag
add_textbox(slide, 4.5, 23.05, 27, 0.5,
            "CMPT 743 · Visual Computing Lab II — Final Project · 2026",
            size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)

# GitHub URL (right)
add_textbox(slide, 31.5, 22.9, 4.0, 0.4,
            "github.com/CMPT743-Team/",
            size=14, color=INK, align=PP_ALIGN.RIGHT)
add_textbox(slide, 31.5, 23.22, 4.0, 0.4,
            "video-scene-text-translator",
            size=14, color=INK, align=PP_ALIGN.RIGHT)

# -- save ----------------------------------------------------------
out = POSTER / "poster_final.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Size: {out.stat().st_size / 1024:.1f} KB")
